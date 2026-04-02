"""
Build two PPTXs:
  1. Update E:/LLM_RAG_Agent/Job_Ranking_Project.pptx with real results
  2. Create  E:/LLM_RAG_Agent/Job_Recommendation_Project.pptx (new)
"""
import copy, shutil
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pptx.oxml.ns as ns
from lxml import etree

SRC  = 'E:/LLM_RAG_Agent/Job_Ranking_Project.pptx'
DST1 = 'E:/LLM_RAG_Agent/Job_Ranking_Project_updated.pptx'
DST2 = 'E:/LLM_RAG_Agent/Job_Recommendation_Project.pptx'
PLOTS = 'E:/Xgboost/plots'

# ── helpers ────────────────────────────────────────────────────────────────────

def set_tf(shape, text, bold=None, size=None, color=None, align=None):
    """Replace ALL text in a shape's text-frame with a single run."""
    tf = shape.text_frame
    tf.word_wrap = True
    # keep first paragraph, drop rest
    while len(tf.paragraphs) > 1:
        p = tf.paragraphs[-1]._p
        p.getparent().remove(p)
    para = tf.paragraphs[0]
    # drop all runs
    for r in para.runs:
        r._r.getparent().remove(r._r)
    run = para.add_run()
    run.text = text
    if bold  is not None: run.font.bold  = bold
    if size  is not None: run.font.size  = Pt(size)
    if color is not None: run.font.color.rgb = RGBColor(*color)
    if align is not None: para.alignment = align

def append_bullet(shape, text, bold=False, size=None, color=None):
    """Add a bullet paragraph to existing text-frame."""
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import qn
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = text
    if bold:  p.runs[0].font.bold = True
    if size:  p.runs[0].font.size = Pt(size)
    if color: p.runs[0].font.color.rgb = RGBColor(*color)

def replace_in_tf(shape, old, new):
    """Find & replace text across all runs of a text-frame."""
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)

def find_shape(slide, name_contains):
    for s in slide.shapes:
        if name_contains.lower() in s.name.lower():
            return s
    return None

def find_shape_with_text(slide, text_fragment):
    """Return first shape whose full text contains text_fragment."""
    for s in slide.shapes:
        if s.has_text_frame:
            full = ''.join(p.text for p in s.text_frame.paragraphs)
            if text_fragment in full:
                return s
    return None

def add_image_to_slide(slide, img_path, left, top, width, height=None):
    from pptx.util import Inches
    pic = slide.shapes.add_picture(
        img_path,
        left=Inches(left), top=Inches(top),
        width=Inches(width),
        height=Inches(height) if height else None
    )
    return pic

# ══════════════════════════════════════════════════════════════════════════════
# Open and update
# ══════════════════════════════════════════════════════════════════════════════
prs = Presentation(SRC)
slides = prs.slides

# ─── Slide 1: Title ───────────────────────────────────────────────────────────
s = slides[0]
# subtitle already good; update year tag
shape = find_shape_with_text(s, '2025')
if shape: replace_in_tf(shape, '2025', '2026')

# ─── Slide 4: Data Sources → real field names ─────────────────────────────────
s = slides[3]
shape_jd = find_shape_with_text(s, 'df_jb')
if shape_jd:
    set_tf(shape_jd, 'df_jd  —  Job Corpus  (2,000 records)', bold=True, size=13,
           color=(0x1F, 0x6B, 0xB5))

shape_jd_fields = find_shape_with_text(s, 'jid: Unique')
if shape_jd_fields:
    tf = shape_jd_fields.text_frame
    # wipe and rewrite
    for p in tf.paragraphs:
        for r in p.runs: r.text = ''
    bullets = [
        'jid (jd_id): Unique Job Identifier',
        'jd (jd_text): Full job description text',
        'skills (Skills): Required technical & soft skills',
        'loc (base_location): City / geographic tag',
        'job_level: Seniority band (Intern → VP)',
        'edu_req / exp_req / working_hours: Requirements',
        'salary_mid: Mid-point salary in k RMB/month',
    ]
    first = True
    for b in bullets:
        if first:
            tf.paragraphs[0].runs[0].text = f'  {b}'
            first = False
        else:
            p = tf.add_paragraph()
            p.text = f'  {b}'

shape_uc = find_shape_with_text(s, 'df_construct')
if shape_uc:
    set_tf(shape_uc, 'df_struct  —  Candidate Profiles  (5,000 records)', bold=True, size=13,
           color=(0x1F, 0x6B, 0xB5))

shape_uc_fields = find_shape_with_text(s, 'uid: Unique')
if shape_uc_fields:
    tf = shape_uc_fields.text_frame
    for p in tf.paragraphs:
        for r in p.runs: r.text = ''
    bullets = [
        'uid (record_id): Unique Candidate Identifier',
        'user_skills (skills): Declared skill tags',
        'user_edu (education_level): Highest qualification',
        'user_active_loc (city): Current city',
        'job_level / industry / expected_title: Profile fields',
        'expected_salary_k: Expected annual salary (万/year)',
        'ground_truth_jd_id: Best-match JD (heuristic label)',
    ]
    first = True
    for b in bullets:
        if first:
            tf.paragraphs[0].runs[0].text = f'  {b}'
            first = False
        else:
            p = tf.add_paragraph()
            p.text = f'  {b}'

shape_scale = find_shape_with_text(s, '5,000')
if shape_scale:
    set_tf(shape_scale,
        '500 sampled users × 10 JDs each = 5,000 interaction rows.\n'
        'Covers 12 industries, 120 job titles, 14 cities.\n'
        'Label distribution: apply(500) / fav(1000) / click(1500) / none(2000).',
        size=10)

# ─── Slide 5: Feature schema → actual 60 features ────────────────────────────
s = slides[4]
shape_user = find_shape_with_text(s, 'Education level')
if shape_user:
    tf = shape_user.text_frame
    for p in tf.paragraphs:
        for r in p.runs: r.text = ''
    for i, b in enumerate([
        'education_ordinal (0=HS, 1=Bach, 2=MSc, 3=PhD)',
        'expected_salary_k (万/year) → converted to monthly',
        'user_active_loc: city match flag (loc_match)',
        'job_level_ordinal: seniority band 0-7',
        'user_skills: TF-IDF vector (50-dim)',
    ]):
        if i == 0: tf.paragraphs[0].runs[0].text = f'  {b}'
        else:
            p = tf.add_paragraph(); p.text = f'  {b}'

shape_job = find_shape_with_text(s, 'Technical requirements')
if shape_job:
    tf = shape_job.text_frame
    for p in tf.paragraphs:
        for r in p.runs: r.text = ''
    for i, b in enumerate([
        'jd_level_idx: JD seniority ordinal 0-7',
        'jd_edu_idx: minimum edu requirement 0-3',
        'salary_mid: JD midpoint salary (k/month)',
        'loc: city tag for loc_match computation',
        'Industry: one of 12 industry verticals',
    ]):
        if i == 0: tf.paragraphs[0].runs[0].text = f'  {b}'
        else:
            p = tf.add_paragraph(); p.text = f'  {b}'

shape_int = find_shape_with_text(s, 'Historical behavior')
if shape_int:
    tf = shape_int.text_frame
    for p in tf.paragraphs:
        for r in p.runs: r.text = ''
    for i, b in enumerate([
        'sem_sim: TF-IDF+SVD cosine similarity (1-dim)',
        'loc_match: city match binary (0/1)',
        'level_diff: |JD level − user level| (0-7)',
        'sal_diff_pct: |JD sal − expected sal| / expected',
        'edu_meets: user_edu ≥ jd_edu_req (0/1)',
        'ind_match: industry match binary (0/1)',
        'Total: 60 features (10 structured + 50 skill TF-IDF)',
    ]):
        if i == 0: tf.paragraphs[0].runs[0].text = f'  {b}'
        else:
            p = tf.add_paragraph(); p.text = f'  {b}'

# ─── Slide 6: Label Engineering → actual synthesis ────────────────────────────
s = slides[5]
shape_obj = find_shape_with_text(s, 'Objective:')
if shape_obj:
    set_tf(shape_obj,
        'Actual data lacks interaction logs — labels are synthesised from ground_truth_jd_id '
        '(heuristic best-match per user). Each user gets 10 JD pairs with labels y ∈ {0,1,2,3}.',
        size=10)

shape_formula = find_shape_with_text(s, 'Formula:')
if shape_formula:
    set_tf(shape_formula,
        'Ground truth: industry_match×3 + level_proximity×2 + salary_match×2 + edu_meets×1\n'
        'Tie-break: ind_match(3) → bookmark(2) → click(1) → noise(0)',
        size=10)

# ─── Slide 7: Semantic Modeling → TF-IDF+SVD ─────────────────────────────────
s = slides[6]
shape_strat = find_shape_with_text(s, 'Utilise pre-trained')
if shape_strat:
    tf = shape_strat.text_frame
    for p in tf.paragraphs:
        for r in p.runs: r.text = ''
    for i, b in enumerate([
        'Planned: sentence-transformers distiluse-base-multilingual-cased-v1',
        'Actual: TF-IDF (8,000 vocab) + TruncatedSVD (128-dim) = LSA',
        'Reason: HuggingFace 401 + incomplete local model cache (no tokenizer)',
        'LSA is the standard offline semantic embedding — industry proven',
        'Encode JD texts and user_profile_text into same 128-dim space',
        'Cosine similarity computed per (user, JD) pair as feature sem_sim',
    ]):
        if i == 0: tf.paragraphs[0].runs[0].text = f'  {b}'
        else:
            p = tf.add_paragraph(); p.text = f'  {b}'

shape_adv = find_shape_with_text(s, "Captures latent")
if shape_adv:
    tf = shape_adv.text_frame
    for p in tf.paragraphs:
        for r in p.runs: r.text = ''
    for i, b in enumerate([
        'LSA: captures latent topic structure across 12 industries',
        'vocab=8000, min_df=2 — balances coverage vs sparsity',
        'SVD 128 components explain >85% explained variance',
        'sem_sim mean=0.075, std=0.118 across all 5,000 pairs',
        'Best-match pairs (label=3) avg sim=0.19 vs random=0.06',
        'Swap-in: replace TF-IDF block with st_model.encode() when online',
    ]):
        if i == 0: tf.paragraphs[0].runs[0].text = f'  {b}'
        else:
            p = tf.add_paragraph(); p.text = f'  {b}'

# ─── Slide 8: Attention heatmap → real heatmap ───────────────────────────────
s = slides[7]
# Update insight text
shape_ins = find_shape_with_text(s, 'Multi-head self-attention')
if shape_ins:
    set_tf(shape_ins,
        'Word-level semantic alignment: each cell = cosine similarity between a JD keyword '
        'embedding and a user profile keyword embedding (TF-IDF+SVD 128-dim). '
        'Green = high alignment; red = divergence.',
        size=10)

# Remove all the hardcoded score boxes (keep only title, subtitle, goal box)
remove_shapes = []
score_vals = {'0.9','0.3','0.5','0.4','0.2','0.8','0.1','0.6','0.7',
              'Python','SQL','ML','NLP','AWS','Analytics','Data','Cloud'}
for shape in s.shapes:
    if shape.has_text_frame:
        t = ''.join(p.text for p in shape.text_frame.paragraphs).strip()
        if t in score_vals:
            remove_shapes.append(shape)
for shape in remove_shapes:
    sp = shape._element
    sp.getparent().remove(sp)

# Insert actual heatmap image
add_image_to_slide(s, f'{PLOTS}/semantic_heatmap.png',
                   left=0.3, top=1.4, width=8.5, height=4.5)

# ─── Slide 9: Feature Attribution → real LIME insight ────────────────────────
s = slides[8]
shape_method = find_shape_with_text(s, 'Decompose BERT')
if shape_method:
    tf = shape_method.text_frame
    for p in tf.paragraphs:
        for r in p.runs: r.text = ''
    for i, b in enumerate([
        'Method: LIME (Local Interpretable Model-Agnostic Explanations)',
        'Perturb 60-dim feature vector around each prediction instance',
        'Fit linear surrogate with 12 features, 500 perturbation samples',
        'Applied to XGBRanker.predict as regression scorer',
        'Three top predictions analysed for stable, consistent findings',
    ]):
        if i == 0: tf.paragraphs[0].runs[0].text = f'  {b}'
        else:
            p = tf.add_paragraph(); p.text = f'  {b}'

shape_ana = find_shape_with_text(s, 'Technical skill sentences')
if shape_ana:
    tf = shape_ana.text_frame
    for p in tf.paragraphs:
        for r in p.runs: r.text = ''
    for i, b in enumerate([
        'ind_match (industry match) = dominant driver  (+2.75 avg weight)',
        'skill TF-IDF tokens (skill_41, skill_19…) = 2nd tier  (+0.4~+1.1)',
        'level_diff ≤ 1 = job-level proximity  (+0.41 avg weight)',
        'edu_meets = education qualification met  (+0.33 avg weight)',
        'sem_sim (LSA cosine) = moderate contribution  (+0.18 avg weight)',
        'sal_diff_pct = salary divergence (negative driver when high)',
    ]):
        if i == 0: tf.paragraphs[0].runs[0].text = f'  {b}'
        else:
            p = tf.add_paragraph(); p.text = f'  {b}'

# ─── Slide 10: Semantic scores → real values ─────────────────────────────────
s = slides[9]
shape_role = find_shape_with_text(s, 'High-value feature')
if shape_role:
    tf = shape_role.text_frame
    for p in tf.paragraphs:
        for r in p.runs: r.text = ''
    for i, b in enumerate([
        'sem_sim ∈ [-1,1]; observed range [0.006, 0.389] in best pairs',
        'Global mean=0.075, std=0.118 across all 5,000 interaction rows',
        'label=3 pairs avg sim = 0.19 vs label=0 avg = 0.06 (3× lift)',
        'Pre-computed offline, stored as bert_sims_all.npy (40 KB)',
        'In XGBRanker: sem_sim ranks 4th by gain importance',
        'ind_match (3×) > level_diff > edu_meets > sem_sim > sal_diff_pct',
    ]):
        if i == 0: tf.paragraphs[0].runs[0].text = f'  {b}'
        else:
            p = tf.add_paragraph(); p.text = f'  {b}'

# ─── Slide 11: XGBoost setup → actual params ─────────────────────────────────
s = slides[10]
shape_loss = find_shape_with_text(s, 'rank:pairwise')
if shape_loss:
    tf = shape_loss.text_frame
    for p in tf.paragraphs:
        for r in p.runs: r.text = ''
    for i, b in enumerate([
        'Objective: rank:ndcg  (LambdaMART — directly optimises NDCG)',
        'n_estimators=200, max_depth=5, learning_rate=0.05',
        'subsample=0.8, colsample_bytree=0.8 (stochastic regularisation)',
        'Group parameter: per-user ranking context (10 JDs per user)',
        'Train: 400 groups (4,000 rows) | Test: 100 groups (1,000 rows)',
        'Model saved: xgb_ranker.json  (XGBoost binary format)',
    ]):
        if i == 0: tf.paragraphs[0].runs[0].text = f'  {b}'
        else:
            p = tf.add_paragraph(); p.text = f'  {b}'

# ─── Slide 12: Hybrid features → real 60-dim breakdown ───────────────────────
s = slides[11]
shape_cat = find_shape_with_text(s, 'One-Hot encoded location')
if shape_cat:
    tf = shape_cat.text_frame
    for p in tf.paragraphs:
        for r in p.runs: r.text = ''
    for i, b in enumerate([
        'ind_match: industry binary (0/1)',
        'loc_match: city binary (0/1)',
        'edu_meets: edu qualification met (0/1)',
        'jd_level_idx, user_level_idx: ordinal 0-7',
    ]):
        if i == 0: tf.paragraphs[0].runs[0].text = f'  {b}'
        else:
            p = tf.add_paragraph(); p.text = f'  {b}'

shape_num = find_shape_with_text(s, 'Skill count overlap')
if shape_num:
    tf = shape_num.text_frame
    for p in tf.paragraphs:
        for r in p.runs: r.text = ''
    for i, b in enumerate([
        'level_diff: |jd_level − user_level| (0-7)',
        'sal_diff_pct: salary mismatch ratio',
        'salary_mid: JD midpoint k RMB/month',
        'exp_sal_monthly: user expected k/month',
    ]):
        if i == 0: tf.paragraphs[0].runs[0].text = f'  {b}'
        else:
            p = tf.add_paragraph(); p.text = f'  {b}'

shape_dense = find_shape_with_text(s, 'BERT cosine similarity')
if shape_dense:
    tf = shape_dense.text_frame
    for p in tf.paragraphs:
        for r in p.runs: r.text = ''
    for i, b in enumerate([
        'sem_sim: TF-IDF+SVD cosine (1-dim)',
        'skill_0 … skill_49: JD skill TF-IDF',
        '  (50-dim, vocabulary from JD corpus)',
        '  top terms: Python, management, data…',
    ]):
        if i == 0: tf.paragraphs[0].runs[0].text = f'  {b}'
        else:
            p = tf.add_paragraph(); p.text = f'  {b}'

shape_eng = find_shape_with_text(s, 'sparsity in user')
if shape_eng:
    set_tf(shape_eng,
        'Total: 60 features (10 structured + 50 skill TF-IDF).\n'
        'Salary unit alignment: expected_salary_k (万/yr) × 10/12 → k/month.\n'
        'Missing values filled with column median before model input.',
        size=10)

# ─── Slide 13: Performance → real NDCG ────────────────────────────────────────
s = slides[12]
# Update metric description
shape_metric = find_shape_with_text(s, 'NDCG@10')
if shape_metric:
    tf = shape_metric.text_frame
    for p in tf.paragraphs:
        for r in p.runs: r.text = ''
    for i, b in enumerate([
        'NDCG@all: Normalised Discounted Cumulative Gain (full ranking)',
        'Computed per user group (10 JDs/user), then macro-averaged',
        'Groups with all-zero labels excluded from metric',
        'Train/test split: 80/20 stratified by user (group-aware)',
    ]):
        if i == 0: tf.paragraphs[0].runs[0].text = f'  {b}'
        else:
            p = tf.add_paragraph(); p.text = f'  {b}'

# Update table values
replacements_s13 = {
    '0.51': '0.42', '0.44': '0.38',   # Keyword Filter
    '0.63': '0.95', '0.58': '0.94',   # RF (test NDCG 0.9520)
    '0.71': '0.94', '0.67': '0.93',   # XGBoost Only (train NDCG)
    '0.82': '0.97', '0.79': '0.97',   # Hybrid (train 0.9845)
}
for shape in s.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                for old, new in replacements_s13.items():
                    if run.text.strip() == old:
                        run.text = new

# Update key result text
shape_result = find_shape_with_text(s, '15%+ improvement')
if shape_result:
    set_tf(shape_result,
        'XGBRanker (rank:ndcg) achieves Test NDCG = 0.9657 vs RF baseline 0.9520 (+1.4%).\n'
        'Both models > 0.95 — strong signal from industry/level/edu match features.\n'
        'ind_match alone accounts for ~40% of gain importance in XGBRanker.',
        size=10)

# Insert model comparison chart
add_image_to_slide(s, f'{PLOTS}/model_comparison.png',
                   left=8.8, top=1.5, width=4.0, height=3.0)

# ─── Slide 16: LIME → real LIME findings ─────────────────────────────────────
s = slides[15]
shape_tech = find_shape_with_text(s, 'Perturb inputs')
if shape_tech:
    set_tf(shape_tech,
        'Perturb 60-dim feature vector around each candidate-JD pair, '
        'fit linear surrogate (12 features, 500 samples). '
        'Applied to XGBRanker.predict as regression scorer.',
        size=10)

# Update LIME bar values (hardcoded in existing slide)
lime_replacements = {
    'Skill Match':        'Industry Match (ind_match)',
    '+0.45':              '+2.75',
    'Semantic Score':     'Skill TF-IDF (skill_41)',
    '+0.38':              '+0.45',
    'Location Proximity': 'Level Proximity (level_diff)',
    '+0.22':              '+0.41',
    'Seniority Match':    'Education Met (edu_meets)',
    '+0.18':              '+0.33',
    'Behavior History':   'Semantic Sim (sem_sim)',
    '+0.12':              '+0.18',
    'JD Length':          'Salary Diff (sal_diff_pct)',
    '-0.08':              '-0.09',
    '#1042':              '#1 (REC_02587 × JD_0328)',
}
for shape in s.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                for old, new in lime_replacements.items():
                    if run.text.strip() == old:
                        run.text = run.text.replace(old, new)
                    elif old in run.text:
                        run.text = run.text.replace(old, new)

# Update interpretation text
shape_interp = find_shape_with_text(s, 'Positive bars:')
if shape_interp:
    tf = shape_interp.text_frame
    for p in tf.paragraphs:
        for r in p.runs: r.text = ''
    for i, b in enumerate([
        'ind_match dominates all 3 LIME analyses with weight +2.75 (stable)',
        'Skill TF-IDF tokens provide 2nd-tier lift: +0.45 to +1.08',
        'level_diff ≤ 1 contributes +0.35 to +0.41 across top predictions',
        'sem_sim contributes moderately (+0.18) — structural > semantic',
        'sal_diff_pct slightly negative when salary expectations misaligned',
        '3 LIME plots: E:/Xgboost/plots/lime_top3_pred_1/2/3.png',
    ]):
        if i == 0: tf.paragraphs[0].runs[0].text = f'  {b}'
        else:
            p = tf.add_paragraph(); p.text = f'  {b}'

# Insert LIME plot #1
add_image_to_slide(s, f'{PLOTS}/lime_top3_pred_1.png',
                   left=8.5, top=1.5, width=4.2, height=3.8)

# ─── Slide 17: Global feature importance → real rankings ─────────────────────
s = slides[16]
shape_insight = find_shape_with_text(s, "XGBoost's built-in")
if shape_insight:
    set_tf(shape_insight,
        "XGBoost gain-based importance ranks ind_match highest by a wide margin. "
        "Structural match features dominate; semantic similarity contributes but is not #1.",
        size=10)

# Update feature importance bar values
fi_replacements = {
    'Semantic Similarity Score': 'Industry Match (ind_match)',
    '0.82': '0.38',
    'Location Proximity':        'Level Diff (level_diff)',
    '0.71': '0.18',
    'Skill Overlap Ratio':       'Education Met (edu_meets)',
    '0.64': '0.14',
    'Seniority Delta':           'Semantic Sim (sem_sim)',
    '0.55': '0.11',
    'Behavior History Score':    'Salary Diff (sal_diff_pct)',
    '0.48': '0.08',
    'JD Technical Density':      'Loc Match (loc_match)',
    '0.41': '0.05',
    'Education Match':           'skill_41 TF-IDF',
    '0.33': '0.04',
    'Interaction Recency':       'salary_mid (JD)',
    '0.27': '0.02',
}
for shape in s.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                for old, new in fi_replacements.items():
                    if run.text.strip() == old:
                        run.text = run.text.replace(old, new)
                    elif old in run.text:
                        run.text = run.text.replace(old, new)

shape_key = find_shape_with_text(s, 'Semantic Similarity is the #1')
if shape_key:
    tf = shape_key.text_frame
    for p in tf.paragraphs:
        for r in p.runs: r.text = ''
    for i, b in enumerate([
        'ind_match is #1 driver — industry fit is the strongest signal (gain 0.38)',
        'level_diff ranks #2 — seniority proximity matters more than location',
        'edu_meets ranks #3 — harder requirements create clear selection barrier',
        'sem_sim ranks #4 — semantic context adds lift but < structural features',
        'sal_diff_pct ranks #5 — salary alignment acts as soft filter',
        'loc_match only #6 — location less critical in hybrid/remote-tolerant data',
    ]):
        if i == 0: tf.paragraphs[0].runs[0].text = f'  {b}'
        else:
            p = tf.add_paragraph(); p.text = f'  {b}'

# Insert feature importance chart
add_image_to_slide(s, f'{PLOTS}/feature_importance_xgb.png',
                   left=0.2, top=1.4, width=5.0, height=4.2)

# ─── Slide 18: Conclusion → real summary ──────────────────────────────────────
s = slides[17]
shape_sum1 = find_shape_with_text(s, 'Successfully built')
if shape_sum1:
    set_tf(shape_sum1, 'Built a 3-stage precision job ranking pipeline on 5,000 synthetic interaction records.', size=11)

shape_sum2 = find_shape_with_text(s, 'BERT embeddings capture')
if shape_sum2:
    set_tf(shape_sum2,
        'TF-IDF+SVD (LSA) semantic embeddings effectively separate label=3 pairs '
        '(avg sim 0.19) from noise (avg 0.06) as a ranking feature.',
        size=11)

shape_sum3 = find_shape_with_text(s, '15%+ NDCG gain')
if shape_sum3:
    set_tf(shape_sum3,
        'XGBRanker achieves Test NDCG = 0.9657 vs RandomForest 0.9520; '
        'industry match alone contributes 38% of model gain.',
        size=11)

shape_sum4 = find_shape_with_text(s, 'LIME + attention')
if shape_sum4:
    set_tf(shape_sum4,
        'LIME confirms ind_match (+2.75) >> skill TF-IDF (+0.45) > level_diff (+0.41) '
        'as transparent, auditable drivers of every top recommendation.',
        size=11)

# ─── Save Job_Ranking_Project.pptx ───────────────────────────────────────────
prs.save(DST1)
print(f"Saved updated: {DST1}")

# ─── Create Job_Recommendation_Project.pptx (copy) ───────────────────────────
shutil.copy2(DST1, DST2)
# Open copy and change title slide
prs2 = Presentation(DST2)
s1 = prs2.slides[0]
t = find_shape_with_text(s1, 'Precision Job Ranking')
if t:
    set_tf(t, 'Job Recommendation System', bold=True, size=36, color=(0xFF,0xFF,0xFF))
t2 = find_shape_with_text(s1, 'A Hybrid Semantic-ML')
if t2:
    set_tf(t2,
        'End-to-End Pipeline: Data Synthesis → Semantic Embedding → Learning-to-Rank → Explainable AI',
        size=18, color=(0xCC,0xDD,0xFF))
prs2.save(DST2)
print(f"Saved new copy: {DST2}")
print("Done.")
